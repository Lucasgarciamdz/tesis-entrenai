from abc import ABC, abstractmethod
from pathlib import Path
from typing import List, Optional, Dict

from entrenai2.configuracion.registrador import obtener_registrador

# Importaciones específicas para procesamiento de archivos
# Se importan aquí para que estén disponibles en el ámbito de las clases y los bloques except
try:
    import pytesseract
    from pdf2image import convert_from_path
    from pdf2image.exceptions import PDFInfoNotInstalledError, PDFPageCountError, PDFSyntaxError
except ImportError:
    pytesseract = None
    convert_from_path = None
    PDFInfoNotInstalledError = type('PDFInfoNotInstalledError', (Exception,), {})
    PDFPageCountError = type('PDFPageCountError', (Exception,), {})
    PDFSyntaxError = type('PDFSyntaxError', (Exception,), {})

try:
    import docx
except ImportError:
    docx = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None


registrador = obtener_registrador(__name__)


class ErrorProcesamientoArchivo(Exception):
    """Excepción personalizada para errores de procesamiento de archivos."""
    pass


class ProcesadorBaseArchivo(ABC):
    """
    Clase base abstracta para procesadores de archivos.
    Cada subclase debe manejar un tipo de archivo específico.
    """
    EXTENSIONES_SOPORTADAS: List[str] = []

    @abstractmethod
    def extraer_texto(self, ruta_archivo: Path) -> str:
        """
        Extrae el contenido textual del archivo dado.
        Debería lanzar ErrorProcesamientoArchivo si la extracción falla.
        """
        pass

    def puede_procesar(self, ruta_archivo: Path) -> bool:
        """
        Verifica si este procesador puede manejar el tipo de archivo dado según su extensión.
        """
        return ruta_archivo.suffix.lower() in self.EXTENSIONES_SOPORTADAS


class ProcesadorArchivoTxt(ProcesadorBaseArchivo):
    """Procesa archivos de texto plano (.txt)."""
    EXTENSIONES_SOPORTADAS = [".txt"]

    def extraer_texto(self, ruta_archivo: Path) -> str:
        registrador.info(f"Extrayendo texto de archivo TXT: {ruta_archivo}")
        codificaciones_a_intentar = [
            "utf-8", "latin-1", "iso-8859-1", "cp1252",
        ]
        
        for codificacion in codificaciones_a_intentar:
            try:
                with open(ruta_archivo, "r", encoding=codificacion) as f:
                    texto = f.read()
                registrador.info(f"Archivo leído exitosamente con codificación: {codificacion}")
                return texto
            except Exception as e:
                registrador.debug(f"Fallo al leer TXT con codificación {codificacion}: {e}")
                continue

        mensaje_error = f"No se pudo extraer texto del archivo TXT {ruta_archivo} con ninguna codificación probada."
        registrador.error(mensaje_error)
        raise ErrorProcesamientoArchivo(mensaje_error)


class ProcesadorArchivoMarkdown(ProcesadorBaseArchivo):
    """Procesa archivos Markdown (.md)."""
    EXTENSIONES_SOPORTADAS = [".md", ".markdown"]

    def extraer_texto(self, ruta_archivo: Path) -> str:
        registrador.info(f"Extrayendo texto de archivo Markdown: {ruta_archivo}")
        try:
            with open(ruta_archivo, "r", encoding="utf-8") as f:
                return f.read()
        except Exception as e:
            registrador.error(f"Error procesando archivo Markdown {ruta_archivo}: {e}")
            raise ErrorProcesamientoArchivo(
                f"No se pudo extraer texto del archivo Markdown {ruta_archivo}: {e}"
            ) from e


class ProcesadorArchivoPdf(ProcesadorBaseArchivo):
    """Procesa archivos PDF (.pdf) usando OCR (Tesseract vía pdf2image)."""
    EXTENSIONES_SOPORTADAS = [".pdf"]

    def extraer_texto(self, ruta_archivo: Path) -> str:
        registrador.info(f"Extrayendo texto de archivo PDF: {ruta_archivo} usando OCR.")
        texto_completo_partes: List[str] = []
        try:
            if pytesseract is None or convert_from_path is None:
                raise ErrorProcesamientoArchivo("Dependencias para PDF (pytesseract, pdf2image) no instaladas.")
            
            pytesseract.get_tesseract_version() # Verifica si Tesseract está instalado
            
            imagenes = convert_from_path(ruta_archivo)
            if not imagenes:
                registrador.warning(
                    f"pdf2image no devolvió imágenes para el PDF: {ruta_archivo}. El PDF podría estar vacío o corrupto."
                )
                return ""

            for i, imagen in enumerate(imagenes):
                registrador.debug(
                    f"Procesando página {i + 1} de {len(imagenes)} del PDF {ruta_archivo} con OCR..."
                )
                try:
                    texto_pagina = pytesseract.image_to_string(imagen, lang="spa+eng")
                    texto_completo_partes.append(texto_pagina)
                except pytesseract.TesseractError as ocr_page_err:
                    registrador.error(
                        f"Error de Tesseract OCR en página {i + 1} de {ruta_archivo}: {ocr_page_err}"
                    )
                    texto_completo_partes.append(
                        f"\n[Error OCR en página {i + 1}: {ocr_page_err}]\n"
                    )
                except Exception as page_processing_err:
                    registrador.error(
                        f"Error inesperado procesando página {i + 1} del PDF {ruta_archivo} con OCR: {page_processing_err}"
                    )
                    texto_completo_partes.append(f"\n[Error procesando página {i + 1}]\n")

            texto_extraido = "\n\n".join(filter(None, texto_completo_partes))
            registrador.info(
                f"Texto extraído exitosamente de PDF (OCR): {ruta_archivo} (longitud: {len(texto_extraido)})"
            )
            return texto_extraido

        except (PDFInfoNotInstalledError, PDFPageCountError, PDFSyntaxError) as pdf_err:
            registrador.error(
                f"Error procesando archivo PDF {ruta_archivo}: {pdf_err}"
            )
            raise ErrorProcesamientoArchivo(
                f"Archivo PDF inválido o corrupto {ruta_archivo}: {pdf_err}"
            ) from pdf_err
        except ErrorProcesamientoArchivo: # Re-raise custom errors
            raise
        except Exception as e:
            registrador.error(
                f"Error general procesando archivo PDF {ruta_archivo} con OCR: {e}"
            )
            raise ErrorProcesamientoArchivo(
                f"No se pudo extraer texto del PDF {ruta_archivo}: {e}"
            ) from e


class ProcesadorArchivoDocx(ProcesadorBaseArchivo):
    """Procesa archivos DOCX (.docx)."""
    EXTENSIONES_SOPORTADAS = [".docx"]

    def extraer_texto(self, ruta_archivo: Path) -> str:
        registrador.info(f"Extrayendo texto de archivo DOCX: {ruta_archivo}")
        try:
            if docx is None:
                raise ErrorProcesamientoArchivo("Dependencia para DOCX (python-docx) no instalada.")
            documento = docx.Document(str(ruta_archivo))
            texto_completo = []
            for parrafo in documento.paragraphs:
                texto_completo.append(parrafo.text)
            for tabla in documento.tables:
                for fila in tabla.rows:
                    for celda in fila.cells:
                        texto_completo.append(celda.text)
            registrador.info(
                f"Texto extraído exitosamente de DOCX: {ruta_archivo} (longitud: {sum(len(t) for t in texto_completo)})"
            )
            return "\n\n".join(texto_completo)
        except ErrorProcesamientoArchivo: # Re-raise custom errors
            raise
        except Exception as e:
            registrador.error(f"Error procesando archivo DOCX {ruta_archivo}: {e}")
            raise ErrorProcesamientoArchivo(
                f"No se pudo extraer texto del DOCX {ruta_archivo}: {e}"
            ) from e


class ProcesadorArchivoPptx(ProcesadorBaseArchivo):
    """Procesa archivos PPTX (.pptx)."""
    EXTENSIONES_SOPORTADAS = [".pptx"]

    def extraer_texto(self, ruta_archivo: Path) -> str:
        registrador.info(f"Extrayendo texto de archivo PPTX: {ruta_archivo}")
        try:
            if Presentation is None:
                raise ErrorProcesamientoArchivo("Dependencia para PPTX (python-pptx) no instalada.")
            presentacion = Presentation(str(ruta_archivo))
            texto_completo = []
            for diapositiva in presentacion.slides:
                for forma in diapositiva.shapes:
                    # Acceder a text_frame de forma segura y verificar su existencia
                    marco_texto = getattr(forma, "text_frame", None)
                    if marco_texto:
                        for parrafo in marco_texto.paragraphs:
                            for ejecucion in parrafo.runs:
                                texto_completo.append(ejecucion.text)
                    # Acceder a text de forma segura y verificar su existencia
                    elif hasattr(forma, "text"): # type: ignore[attr-defined]
                        texto_forma = getattr(forma, "text", None)
                        if texto_forma:
                            texto_completo.append(texto_forma)
                
                # Acceder a notes_slide y notes_text_frame de forma segura
                if diapositiva.has_notes_slide:
                    diapositiva_notas = getattr(diapositiva, "notes_slide", None)
                    if diapositiva_notas:
                        marco_texto_notas = getattr(diapositiva_notas, "notes_text_frame", None)
                        if marco_texto_notas:
                            texto_completo.append(marco_texto_notas.text)

            registrador.info(
                f"Texto extraído exitosamente de PPTX: {ruta_archivo} (longitud: {sum(len(t) for t in texto_completo)})"
            )
            return "\n\n".join(filter(None, texto_completo))
        except ErrorProcesamientoArchivo: # Re-raise custom errors
            raise
        except Exception as e:
            registrador.error(f"Error procesando archivo PPTX {ruta_archivo}: {e}")
            raise ErrorProcesamientoArchivo(
                f"No se pudo extraer texto del PPTX {ruta_archivo}: {e}"
            ) from e


class GestorProcesadoresArchivos:
    def __init__(self):
        self.procesadores: Dict[str, ProcesadorBaseArchivo] = {}
        self._registrar_procesadores_por_defecto()

    def _registrar_procesadores_por_defecto(self):
        self.registrar_procesador(ProcesadorArchivoTxt())
        self.registrar_procesador(ProcesadorArchivoMarkdown())
        self.registrar_procesador(ProcesadorArchivoPdf())
        self.registrar_procesador(ProcesadorArchivoDocx())
        self.registrar_procesador(ProcesadorArchivoPptx())

    def registrar_procesador(self, procesador: ProcesadorBaseArchivo):
        extensiones_soportadas = getattr(procesador, "EXTENSIONES_SOPORTADAS", [])
        if not extensiones_soportadas:
            registrador.warning(
                f"Procesador {type(procesador).__name__} no tiene EXTENSIONES_SOPORTADAS. No se puede registrar."
            )
            return

        for ext in extensiones_soportadas:
            ext_minusculas = ext.lower()
            if ext_minusculas not in self.procesadores:
                self.procesadores[ext_minusculas] = procesador
                registrador.info(
                    f"Registrado {type(procesador).__name__} para extensión {ext_minusculas}"
                )
            else:
                registrador.warning(
                    f"Procesador para extensión {ext_minusculas} ya registrado con {type(self.procesadores[ext_minusculas]).__name__}. "
                    f"Omitiendo registro de {type(procesador).__name__}."
                )

    def procesar_archivo(self, ruta_archivo: Path) -> Optional[str]:
        if not ruta_archivo.is_file():
            registrador.error(f"Archivo no encontrado o no es un archivo: {ruta_archivo}")
            return None

        extension_archivo = ruta_archivo.suffix.lower()
        procesador = self.procesadores.get(extension_archivo)

        if procesador:
            try:
                registrador.info(
                    f"Procesando archivo '{ruta_archivo}' usando {type(procesador).__name__}..."
                )
                return procesador.extraer_texto(ruta_archivo)
            except ErrorProcesamientoArchivo as e:
                registrador.error(f"Procesamiento de archivo falló para '{ruta_archivo}': {e}")
                return None
            except Exception as e:
                registrador.exception(
                    f"Error inesperado durante procesamiento de archivo para '{ruta_archivo}': {e}"
                )
                return None
        else:
            registrador.warning(
                f"No hay procesador registrado para el tipo de archivo '{extension_archivo}' (archivo: {ruta_archivo}). Omitiendo."
            )
            return None
