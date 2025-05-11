from abc import ABC, abstractmethod
from pathlib import Path
from typing import List, Optional, Dict

from src.entrenai.utils.logger import get_logger

# Import specific libraries for file processing
# For PDF:
from pdf2image import (
    convert_from_path,
    PDFInfoNotInstalledError,
    PDFPageCountError,
    PDFSyntaxError,
)  # type: ignore
import pytesseract  # type: ignore

# For DOCX:
import docx  # type: ignore

# For PPTX:
from pptx import Presentation  # type: ignore

logger = get_logger(__name__)


class FileProcessingError(Exception):
    """Custom exception for file processing errors."""

    pass


class BaseFileProcessor(ABC):
    """
    Abstract base class for file processors.
    Each subclass should handle a specific file type.
    """

    SUPPORTED_EXTENSIONS: List[str] = []

    @abstractmethod
    def extract_text(self, file_path: Path) -> str:
        """
        Extracts text content from the given file.
        Should raise FileProcessingError if extraction fails.
        """
        pass

    def can_process(self, file_path: Path) -> bool:
        """
        Checks if this processor can handle the given file type based on extension.
        """
        return file_path.suffix.lower() in self.SUPPORTED_EXTENSIONS


class TxtFileProcessor(BaseFileProcessor):
    """Processes plain text (.txt) files."""

    SUPPORTED_EXTENSIONS = [".txt"]

    def extract_text(self, file_path: Path) -> str:
        logger.info(f"Extracting text from TXT file: {file_path}")
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                return f.read()
        except Exception as e:
            logger.error(f"Error processing TXT file {file_path}: {e}")
            raise FileProcessingError(
                f"Failed to extract text from TXT file {file_path}: {e}"
            ) from e


class MarkdownFileProcessor(BaseFileProcessor):
    """Processes Markdown (.md) files."""

    SUPPORTED_EXTENSIONS = [".md", ".markdown"]

    def extract_text(self, file_path: Path) -> str:
        logger.info(f"Extracting text from Markdown file: {file_path}")
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                return f.read()
        except Exception as e:
            logger.error(f"Error processing Markdown file {file_path}: {e}")
            raise FileProcessingError(
                f"Failed to extract text from Markdown file {file_path}: {e}"
            ) from e


class PdfFileProcessor(BaseFileProcessor):
    """Processes PDF (.pdf) files using OCR (Tesseract via pdf2image)."""

    SUPPORTED_EXTENSIONS = [".pdf"]

    def extract_text(self, file_path: Path) -> str:
        logger.info(f"Extracting text from PDF file: {file_path} using OCR.")
        full_text_parts: List[str] = []
        try:
            # Attempt to get Tesseract version to check if it's installed and accessible.
            # This can raise EnvironmentError (FileNotFoundError on Windows, TesseractNotFoundError on Linux/macOS)
            try:
                pytesseract.get_tesseract_version()
            except Exception as tess_err:  # Catching broad Exception as specific error varies by OS/setup
                logger.error(
                    f"Tesseract OCR is not installed or not found in PATH: {tess_err}"
                )
                raise FileProcessingError(
                    "Tesseract OCR is not installed or not found in PATH."
                ) from tess_err

            images = convert_from_path(
                file_path
            )  # Defaults: dpi=200, fmt='ppm', thread_count=1
            if not images:
                logger.warning(
                    f"pdf2image returned no images for PDF: {file_path}. The PDF might be empty or corrupted."
                )
                return ""

            for i, image in enumerate(images):
                logger.debug(
                    f"Processing page {i + 1} of {len(images)} from PDF {file_path} with OCR..."
                )
                try:
                    # Configure Tesseract language if needed, e.g., lang='eng+spa' for English and Spanish
                    # Default is 'eng'. User should ensure language packs are installed for Tesseract.
                    page_text = pytesseract.image_to_string(
                        image, lang="spa+eng"
                    )  # Example: Spanish + English
                    full_text_parts.append(page_text)
                except pytesseract.TesseractError as ocr_page_err:
                    logger.error(
                        f"Tesseract OCR error on page {i + 1} of {file_path}: {ocr_page_err}"
                    )
                    full_text_parts.append(
                        f"\n[OCR Error on page {i + 1}: {ocr_page_err}]\n"
                    )  # Add error marker
                except Exception as page_processing_err:
                    logger.error(
                        f"Unexpected error processing page {i + 1} of PDF {file_path} with OCR: {page_processing_err}"
                    )
                    full_text_parts.append(f"\n[Error processing page {i + 1}]\n")

            extracted_text = "\n\n".join(
                filter(None, full_text_parts)
            )  # Join non-empty text parts
            logger.info(
                f"Successfully extracted text from PDF (OCR): {file_path} (length: {len(extracted_text)})"
            )
            return extracted_text

        except PDFInfoNotInstalledError as poppler_err:
            logger.error(
                f"Poppler (pdfinfo) not found. pdf2image cannot process PDF files: {poppler_err}"
            )
            raise FileProcessingError(
                "Poppler (dependency for PDF processing) not installed or not in PATH."
            ) from poppler_err
        except (PDFPageCountError, PDFSyntaxError) as pdf_err:
            logger.error(
                f"Error processing PDF file structure for {file_path}: {pdf_err}"
            )
            raise FileProcessingError(
                f"Invalid or corrupted PDF file {file_path}: {pdf_err}"
            ) from pdf_err
        except Exception as e:
            logger.error(f"General error processing PDF file {file_path} with OCR: {e}")
            # Check if it's a TesseractNotFoundError (if not caught by specific TesseractError above)
            # This check might be redundant if the initial get_tesseract_version() check is robust.
            if "Tesseract is not installed or not in your PATH" in str(e):
                raise FileProcessingError(
                    "Tesseract (for OCR) is not installed or not in PATH."
                ) from e
            raise FileProcessingError(
                f"Failed to extract text from PDF {file_path}: {e}"
            ) from e


class DocxFileProcessor(BaseFileProcessor):
    """Processes DOCX (.docx) files."""

    SUPPORTED_EXTENSIONS = [".docx"]

    def extract_text(self, file_path: Path) -> str:
        logger.info(f"Extracting text from DOCX file: {file_path}")
        try:
            document = docx.Document(str(file_path))  # Convert Path to str
            full_text = []
            for para in document.paragraphs:
                full_text.append(para.text)
            for table in document.tables:
                for row in table.rows:
                    for cell in row.cells:
                        full_text.append(cell.text)
            logger.info(
                f"Successfully extracted text from DOCX: {file_path} (length: {sum(len(t) for t in full_text)})"
            )
            return "\n\n".join(full_text)
        except Exception as e:
            logger.error(f"Error processing DOCX file {file_path}: {e}")
            raise FileProcessingError(
                f"Failed to extract text from DOCX {file_path}: {e}"
            ) from e


class PptxFileProcessor(BaseFileProcessor):
    """Processes PPTX (.pptx) files."""

    SUPPORTED_EXTENSIONS = [".pptx"]

    def extract_text(self, file_path: Path) -> str:
        logger.info(f"Extracting text from PPTX file: {file_path}")
        try:
            prs = Presentation(str(file_path))  # Convert Path to str
            full_text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame") and shape.text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                full_text.append(run.text)
                    elif hasattr(shape, "text"):
                        full_text.append(shape.text)
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    full_text.append(slide.notes_slide.notes_text_frame.text)

            logger.info(
                f"Successfully extracted text from PPTX: {file_path} (length: {sum(len(t) for t in full_text)})"
            )
            return "\n\n".join(filter(None, full_text))
        except Exception as e:
            logger.error(f"Error processing PPTX file {file_path}: {e}")
            raise FileProcessingError(
                f"Failed to extract text from PPTX {file_path}: {e}"
            ) from e


class FileProcessor:
    def __init__(self):
        self.processors: Dict[str, BaseFileProcessor] = {}
        self._register_default_processors()

    def _register_default_processors(self):
        self.register_processor(TxtFileProcessor())
        self.register_processor(MarkdownFileProcessor())
        self.register_processor(PdfFileProcessor())
        self.register_processor(DocxFileProcessor())
        self.register_processor(PptxFileProcessor())

    def register_processor(self, processor: BaseFileProcessor):
        supported_extensions = getattr(processor, "SUPPORTED_EXTENSIONS", [])
        if not supported_extensions:
            logger.warning(
                f"Processor {type(processor).__name__} has no SUPPORTED_EXTENSIONS. Cannot register."
            )
            return

        for ext in supported_extensions:
            ext_lower = ext.lower()
            if ext_lower not in self.processors:
                self.processors[ext_lower] = processor
                logger.info(
                    f"Registered {type(processor).__name__} for extension {ext_lower}"
                )
            else:
                logger.warning(
                    f"Processor for extension {ext_lower} already registered with {type(self.processors[ext_lower]).__name__}. "
                    f"Skipping registration of {type(processor).__name__}."
                )

    def process_file(self, file_path: Path) -> Optional[str]:
        if not file_path.is_file():
            logger.error(f"File not found or is not a file: {file_path}")
            return None

        file_ext = file_path.suffix.lower()
        processor = self.processors.get(file_ext)

        if processor:
            try:
                logger.info(
                    f"Processing file '{file_path}' using {type(processor).__name__}..."
                )
                return processor.extract_text(file_path)
            except FileProcessingError as e:
                logger.error(f"File processing failed for '{file_path}': {e}")
                return None
            except Exception as e:
                logger.exception(
                    f"Unexpected error during file processing for '{file_path}': {e}"
                )
                return None
        else:
            logger.warning(
                f"No processor registered for file type '{file_ext}' (file: {file_path}). Skipping."
            )
            return None


if __name__ == "__main__":
    test_dir = Path("test_files_temp_file_processor")
    test_dir.mkdir(exist_ok=True)

    (test_dir / "sample.txt").write_text(
        "This is a simple text file.\nHello world from TXT."
    )
    (test_dir / "sample.md").write_text(
        "# Markdown Header\n\nThis is *markdown* from MD file."
    )

    try:
        doc = docx.Document()
        doc.add_paragraph("This is a paragraph in a DOCX file.")
        doc.add_paragraph("Another paragraph with some more text.")
        doc.add_heading("Heading 1 in DOCX", level=1)
        table = doc.add_table(rows=2, cols=2)
        table.cell(0, 0).text = "Foo"
        table.cell(0, 1).text = "Bar"
        table.cell(1, 0).text = "Baz"
        table.cell(1, 1).text = "Qux"
        doc.save(str(test_dir / "sample.docx"))
        logger.info(f"Created dummy DOCX file: {test_dir / 'sample.docx'}")
    except Exception as e:
        logger.error(
            f"Could not create dummy DOCX file for testing (python-docx might not be installed or error): {e}"
        )

    try:
        prs = Presentation()
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        if title:
            title.text = "Hello from PPTX Title"

        slide_w_emu = prs.slide_width
        slide_h_emu = prs.slide_height

        if slide_w_emu is not None and slide_h_emu is not None:
            left = slide_w_emu // 4
            top = slide_h_emu // 4
            width = slide_w_emu // 2
            height = slide_h_emu // 2

            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.text = "This is text in a textbox on a slide."
            p = tf.add_paragraph()
            p.text = "Another paragraph in the same textbox."
        else:
            logger.warning(
                "Could not determine slide dimensions for dummy PPTX, skipping textbox creation."
            )

        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        if text_frame:
            text_frame.text = "This is a note on the slide."
        prs.save(str(test_dir / "sample.pptx"))
        logger.info(f"Created dummy PPTX file: {test_dir / 'sample.pptx'}")
    except Exception as e:
        logger.error(
            f"Could not create dummy PPTX file for testing (python-pptx might not be installed or error): {e}"
        )

    # Create a dummy PDF for testing (content will not be read by touch, real PDF needed for OCR test)
    (test_dir / "sample.pdf").write_text(
        "%PDF-1.4\n%âãÏÓ\n1 0 obj<</Type/Catalog/Pages 2 0 R>>endobj\n2 0 obj<</Type/Pages/Count 1/Kids[3 0 R]>>endobj\n3 0 obj<</Type/Page/MediaBox[0 0 612 792]/Parent 2 0 R/Resources<<>>>>endobj\nxref\n0 4\n0000000000 65535 f\n0000000010 00000 n\n0000000059 00000 n\n0000000118 00000 n\ntrailer<</Size 4/Root 1 0 R>>\nstartxref\n196\n%%EOF"
    )
    logger.info(f"Created dummy PDF file: {test_dir / 'sample.pdf'}")

    (test_dir / "unsupported.xyz").touch()

    file_processor = FileProcessor()

    for test_file in sorted(test_dir.iterdir()):
        if test_file.is_file():
            print(f"\n--- Processing {test_file.name} ---")
            extracted_content = file_processor.process_file(test_file)
            if extracted_content is not None:
                print(
                    f"Extracted Content (first 200 chars): {extracted_content[:200].replace(chr(10), ' ')}..."
                )
            else:
                print("Could not extract content or file type not supported.")

    import shutil

    try:
        shutil.rmtree(test_dir)
        print(f"\nCleaned up test directory: {test_dir}")
    except Exception as e:
        print(f"Error cleaning up test directory {test_dir}: {e}")
