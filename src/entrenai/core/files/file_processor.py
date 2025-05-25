from abc import ABC, abstractmethod
from pathlib import Path
from typing import List, Optional, Dict

# For DOCX:
import docx  # type: ignore
import pytesseract  # type: ignore

# Import specific libraries for file processing
# For PDF:
from pdf2image import convert_from_path
from pdf2image.exceptions import (
    PDFInfoNotInstalledError,
    PDFPageCountError,
    PDFSyntaxError,
)  # type: ignore

# For PPTX:
from pptx import Presentation  # type: ignore

from src.entrenai.config.logger import get_logger

logger = get_logger(__name__)


class FileProcessingError(Exception):
    """Custom exception for file processing errors."""

    pass


class BaseFileProcessor(ABC):
    """
    Abstract base class for file processors.
    Each subclass should handle a specific file type.
    """

    SUPPORTED_EXTENSIONS: List[str] = []

    @abstractmethod
    def extract_text(self, file_path: Path) -> str:
        """
        Extrae el contenido textual del archivo dado.
        Debería lanzar FileProcessingError si la extracción falla.
        """
        pass

    def can_process(self, file_path: Path) -> bool:
        """
        Verifica si este procesador puede manejar el tipo de archivo dado según su extensión.
        """
        return file_path.suffix.lower() in self.SUPPORTED_EXTENSIONS


class TxtFileProcessor(BaseFileProcessor):
    """Processes plain text (.txt) files."""

    SUPPORTED_EXTENSIONS = [".txt"]

    def extract_text(self, file_path: Path) -> str:
        logger.info(f"Extrayendo texto de archivo TXT: {file_path}")
        encodings_to_try = [
            "utf-8",
            "latin-1",
            "iso-8859-1",
            "cp1252",
        ]  # Common encodings
        errors = []

        for encoding in encodings_to_try:
            try:
                with open(file_path, "r", encoding=encoding) as f:
                    text = f.read()
                logger.info(f"Archivo leído exitosamente con codificación: {encoding}")
                return text
            except Exception as e:
                errors.append(f"Fallo con codificación {encoding}: {e}")
                continue

        error_msg = f"No se pudo extraer texto del archivo TXT {file_path} con ninguna codificación. Errores: {', '.join(errors)}"
        logger.error(error_msg)
        raise FileProcessingError(error_msg)


class MarkdownFileProcessor(BaseFileProcessor):
    """Procesa archivos Markdown (.md)."""

    SUPPORTED_EXTENSIONS = [".md", ".markdown"]

    def extract_text(self, file_path: Path) -> str:
        logger.info(f"Extrayendo texto de archivo Markdown: {file_path}")
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                return f.read()
        except Exception as e:
            logger.error(f"Error procesando archivo Markdown {file_path}: {e}")
            raise FileProcessingError(
                f"No se pudo extraer texto del archivo Markdown {file_path}: {e}"
            ) from e


class PdfFileProcessor(BaseFileProcessor):
    """Procesa archivos PDF (.pdf) usando OCR (Tesseract vía pdf2image)."""

    SUPPORTED_EXTENSIONS = [".pdf"]

    def extract_text(self, file_path: Path) -> str:
        logger.info(f"Extrayendo texto de archivo PDF: {file_path} usando OCR.")
        full_text_parts: List[str] = []
        try:
            try:
                pytesseract.get_tesseract_version()
            except Exception as tess_err:
                logger.error(
                    f"Tesseract OCR no está instalado o no se encuentra en el PATH: {tess_err}"
                )
                raise FileProcessingError(
                    "Tesseract OCR no está instalado o no se encuentra en el PATH."
                ) from tess_err

            images = convert_from_path(file_path)
            if not images:
                logger.warning(
                    f"pdf2image no devolvió imágenes para el PDF: {file_path}. El PDF podría estar vacío o corrupto."
                )
                return ""

            for i, image in enumerate(images):
                logger.debug(
                    f"Procesando página {i + 1} de {len(images)} del PDF {file_path} con OCR..."
                )
                try:
                    page_text = pytesseract.image_to_string(image, lang="spa+eng")
                    full_text_parts.append(page_text)
                except pytesseract.TesseractError as ocr_page_err:
                    logger.error(
                        f"Error de Tesseract OCR en página {i + 1} de {file_path}: {ocr_page_err}"
                    )
                    full_text_parts.append(
                        f"\n[Error OCR en página {i + 1}: {ocr_page_err}]\n"
                    )
                except Exception as page_processing_err:
                    logger.error(
                        f"Error inesperado procesando página {i + 1} del PDF {file_path} con OCR: {page_processing_err}"
                    )
                    full_text_parts.append(f"\n[Error procesando página {i + 1}]\n")

            extracted_text = "\n\n".join(filter(None, full_text_parts))
            logger.info(
                f"Texto extraído exitosamente de PDF (OCR): {file_path} (longitud: {len(extracted_text)})"
            )
            return extracted_text

        except PDFInfoNotInstalledError as poppler_err:
            logger.error(
                f"Poppler (pdfinfo) no encontrado. pdf2image no puede procesar archivos PDF: {poppler_err}"
            )
            raise FileProcessingError(
                "Poppler (dependencia para procesamiento de PDF) no instalado o no en PATH."
            ) from poppler_err
        except (PDFPageCountError, PDFSyntaxError) as pdf_err:
            logger.error(
                f"Error procesando estructura de archivo PDF para {file_path}: {pdf_err}"
            )
            raise FileProcessingError(
                f"Archivo PDF inválido o corrupto {file_path}: {pdf_err}"
            ) from pdf_err
        except Exception as e:
            logger.error(
                f"Error general procesando archivo PDF {file_path} con OCR: {e}"
            )
            if "Tesseract is not installed or not in your PATH" in str(e):
                raise FileProcessingError(
                    "Tesseract (para OCR) no está instalado o no se encuentra en el PATH."
                ) from e
            raise FileProcessingError(
                f"No se pudo extraer texto del PDF {file_path}: {e}"
            ) from e


class DocxFileProcessor(BaseFileProcessor):
    """Procesa archivos DOCX (.docx)."""

    SUPPORTED_EXTENSIONS = [".docx"]

    def extract_text(self, file_path: Path) -> str:
        logger.info(f"Extrayendo texto de archivo DOCX: {file_path}")
        try:
            document = docx.Document(str(file_path))
            full_text = []
            for para in document.paragraphs:
                full_text.append(para.text)
            for table in document.tables:
                for row in table.rows:
                    for cell in row.cells:
                        full_text.append(cell.text)
            logger.info(
                f"Texto extraído exitosamente de DOCX: {file_path} (longitud: {sum(len(t) for t in full_text)})"
            )
            return "\n\n".join(full_text)
        except Exception as e:
            logger.error(f"Error procesando archivo DOCX {file_path}: {e}")
            raise FileProcessingError(
                f"No se pudo extraer texto del DOCX {file_path}: {e}"
            ) from e


class PptxFileProcessor(BaseFileProcessor):
    """Procesa archivos PPTX (.pptx)."""

    SUPPORTED_EXTENSIONS = [".pptx"]

    def extract_text(self, file_path: Path) -> str:
        logger.info(f"Extrayendo texto de archivo PPTX: {file_path}")
        try:
            prs = Presentation(str(file_path))
            full_text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame") and shape.text_frame:
                        # Pylance might still complain here as hasattr doesn't change type
                        text_frame_obj = getattr(
                            shape, "text_frame", None
                        )  # Use getattr for safer access
                        if text_frame_obj:
                            for paragraph in text_frame_obj.paragraphs:  # type: ignore[attr-defined]
                                for run in paragraph.runs:
                                    full_text.append(run.text)
                    elif hasattr(shape, "text"):
                        shape_text = getattr(shape, "text", None)  # Use getattr
                        if shape_text:
                            full_text.append(shape_text)  # type: ignore[arg-type]
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes_text_frame_obj = getattr(
                        slide.notes_slide, "notes_text_frame", None
                    )
                    if notes_text_frame_obj:
                        full_text.append(notes_text_frame_obj.text)  # type: ignore[attr-defined]

            logger.info(
                f"Texto extraído exitosamente de PPTX: {file_path} (longitud: {sum(len(t) for t in full_text)})"
            )
            return "\n\n".join(filter(None, full_text))
        except Exception as e:
            logger.error(f"Error procesando archivo PPTX {file_path}: {e}")
            raise FileProcessingError(
                f"No se pudo extraer texto del PPTX {file_path}: {e}"
            ) from e


class FileProcessor:
    def __init__(self):
        self.processors: Dict[str, BaseFileProcessor] = {}
        self._register_default_processors()

    def _register_default_processors(self):
        self.register_processor(TxtFileProcessor())
        self.register_processor(MarkdownFileProcessor())
        self.register_processor(PdfFileProcessor())
        self.register_processor(DocxFileProcessor())
        self.register_processor(PptxFileProcessor())

    def register_processor(self, processor: BaseFileProcessor):
        supported_extensions = getattr(processor, "SUPPORTED_EXTENSIONS", [])
        if not supported_extensions:
            logger.warning(
                f"Procesador {type(processor).__name__} no tiene SUPPORTED_EXTENSIONS. No se puede registrar."
            )
            return

        for ext in supported_extensions:
            ext_lower = ext.lower()
            if ext_lower not in self.processors:
                self.processors[ext_lower] = processor
                logger.info(
                    f"Registrado {type(processor).__name__} para extensión {ext_lower}"
                )
            else:
                logger.warning(
                    f"Procesador para extensión {ext_lower} ya registrado con {type(self.processors[ext_lower]).__name__}. "
                    f"Omitiendo registro de {type(processor).__name__}."
                )

    def process_file(self, file_path: Path) -> Optional[str]:
        if not file_path.is_file():
            logger.error(f"Archivo no encontrado o no es un archivo: {file_path}")
            return None

        file_ext = file_path.suffix.lower()
        processor = self.processors.get(file_ext)

        if processor:
            try:
                logger.info(
                    f"Procesando archivo '{file_path}' usando {type(processor).__name__}..."
                )
                return processor.extract_text(file_path)
            except FileProcessingError as e:
                logger.error(f"Procesamiento de archivo falló para '{file_path}': {e}")
                return None
            except Exception as e:
                logger.exception(
                    f"Error inesperado durante procesamiento de archivo para '{file_path}': {e}"
                )
                return None
        else:
            logger.warning(
                f"No hay procesador registrado para el tipo de archivo '{file_ext}' (archivo: {file_path}). Omitiendo."
            )
            return None
