from pathlib import Path
from typing import Optional, Dict, List, Any

# --- Importaciones opcionales de bibliotecas de terceros ---
# Estas dependencias deben estar listadas en el archivo requirements.txt del proyecto.
try:
    import docx # Biblioteca para leer archivos .docx (Microsoft Word)
except ImportError:
    docx = None # Se asigna None si la biblioteca no está instalada

try:
    import pytesseract # Biblioteca para OCR (Reconocimiento Óptico de Caracteres)
    from pdf2image import convert_from_path # Utilidad para convertir páginas de PDF a imágenes
except ImportError:
    pytesseract = None
    convert_from_path = None

try:
    from pptx import Presentation # Biblioteca para leer archivos .pptx (Microsoft PowerPoint)
except ImportError:
    Presentation = None

from entrenai_refactor.config.registrador import obtener_registrador

registrador = obtener_registrador(__name__) # Registrador específico para este módulo

# --- Definiciones de Excepciones Personalizadas ---

class ErrorProcesamientoArchivo(Exception):
    """Excepción base para errores ocurridos durante el procesamiento de archivos."""
    def __init__(self, mensaje: str, error_original: Optional[Exception] = None, ruta_archivo: Optional[Path] = None):
        super().__init__(mensaje)
        self.error_original = error_original
        self.ruta_archivo = ruta_archivo
        detalle_archivo = f", Archivo: {ruta_archivo}" if ruta_archivo else ""
        registrador.debug(f"Excepción ErrorProcesamientoArchivo creada: '{mensaje}'{detalle_archivo}, Original: {error_original}")

    def __str__(self):
        detalle_archivo = f" (Archivo: {self.ruta_archivo})" if self.ruta_archivo else ""
        if self.error_original:
            return f"{super().__str__()}{detalle_archivo} (Error original: {type(self.error_original).__name__}: {str(self.error_original)})"
        return f"{super().__str__()}{detalle_archivo}"

class ErrorTipoArchivoNoSoportado(ErrorProcesamientoArchivo):
    """Excepción para cuando se intenta procesar un tipo de archivo no soportado por el sistema."""
    pass

class ErrorDependenciaFaltante(ErrorProcesamientoArchivo):
    """Excepción para cuando falta una dependencia de software necesaria para procesar un tipo de archivo."""
    pass

# --- Interfaz Base para Procesadores de Archivos ---

class ProcesadorArchivoInterfaz:
    """
    Clase base abstracta (interfaz informal) que define la estructura esperada
    para los procesadores de archivos específicos para cada tipo de documento.
    """

    # Lista de extensiones de archivo (en minúsculas, comenzando con punto) que este procesador maneja.
    # Ejemplo: [".txt", ".log"]
    EXTENSIONES_ARCHIVOS_SOPORTADAS: List[str] = []

    def extraer_texto_de_archivo(self, ruta_archivo_entrada: Path) -> str: # Parámetro renombrado
        """
        Método abstracto para extraer contenido textual de un archivo.
        Debe ser implementado obligatoriamente por todas las subclases concretas.

        Args:
            ruta_archivo_entrada: Objeto Path apuntando al archivo a procesar.

        Returns:
            Un string con el texto extraído del archivo.

        Raises:
            NotImplementedError: Si la subclase no implementa este método.
            ErrorProcesamientoArchivo: Si ocurre un error durante la extracción.
        """
        nombre_clase_actual = self.__class__.__name__
        mensaje_error_no_implementado = f"El método 'extraer_texto_de_archivo' no ha sido implementado en la clase '{nombre_clase_actual}'."
        registrador.error(mensaje_error_no_implementado)
        raise NotImplementedError(mensaje_error_no_implementado)

    def puede_procesar_extension(self, ruta_archivo_entrada: Path) -> bool: # Parámetro renombrado
        """
        Verifica si este procesador es capaz de manejar la extensión del archivo proporcionado.
        La comparación es insensible a mayúsculas/minúsculas.

        Args:
            ruta_archivo_entrada: Objeto Path apuntando al archivo.

        Returns:
            True si el procesador soporta la extensión, False en caso contrario.
        """
        extension_archivo_actual = ruta_archivo_entrada.suffix.lower() # Obtener extensión y convertir a minúsculas
        puede_procesar = extension_archivo_actual in self.EXTENSIONES_ARCHIVOS_SOPORTADAS
        registrador.debug(f"Procesador '{self.__class__.__name__}': ¿Puede procesar extensión '{extension_archivo_actual}'? {'Sí' if puede_procesar else 'No'}.")
        return puede_procesar

# --- Implementaciones de Procesadores Específicos por Tipo de Archivo ---

class ProcesadorArchivosTextoPlano(ProcesadorArchivoInterfaz):
    """Procesador especializado para archivos de texto plano (ej. .txt, .log, .csv)."""
    EXTENSIONES_ARCHIVOS_SOPORTADAS = [".txt", ".text", ".log", ".csv", ".tsv"] # Ampliada lista de extensiones comunes

    def extraer_texto_de_archivo(self, ruta_archivo_entrada: Path) -> str:
        registrador.info(f"Intentando extraer texto del archivo de texto plano: '{ruta_archivo_entrada}'.")
        # Lista de codificaciones comunes a intentar, en orden de probabilidad o preferencia.
        codificaciones_comunes_a_intentar = ["utf-8", "latin-1", "iso-8859-1", "cp1252"]

        for codificacion_prueba in codificaciones_comunes_a_intentar:
            try:
                with open(ruta_archivo_entrada, "r", encoding=codificacion_prueba) as archivo_abierto:
                    texto_extraido_del_archivo = archivo_abierto.read()
                registrador.info(f"Texto extraído de '{ruta_archivo_entrada}' utilizando la codificación '{codificacion_prueba}'.")
                return texto_extraido_del_archivo
            except UnicodeDecodeError: # Error específico al intentar decodificar con una codificación incorrecta
                registrador.debug(f"Falló la decodificación del archivo '{ruta_archivo_entrada}' con la codificación '{codificacion_prueba}'. Intentando siguiente codificación.")
                continue # Probar la siguiente codificación en la lista
            except IOError as e_error_io: # Errores de lectura/escritura del archivo
                mensaje_error_io_especifico = f"Error de E/S al leer el archivo de texto '{ruta_archivo_entrada}' con codificación '{codificacion_prueba}': {e_error_io}"
                registrador.error(mensaje_error_io_especifico)
                raise ErrorProcesamientoArchivo(mensaje_error_io_especifico, e_error_io, ruta_archivo=ruta_archivo_entrada) from e_error_io
            except Exception as e_error_inesperado: # Otros errores no previstos
                registrador.warning(f"Error inesperado leyendo '{ruta_archivo_entrada}' con '{codificacion_prueba}': {e_error_inesperado}. Se intentará con otra codificación si es posible.")
                continue # Continuar con la siguiente codificación por si acaso

        # Si todas las codificaciones fallan
        mensaje_error_extraccion_final = f"No se pudo extraer texto del archivo '{ruta_archivo_entrada}' después de intentar con las codificaciones: {', '.join(codificaciones_comunes_a_intentar)}."
        registrador.error(mensaje_error_extraccion_final)
        raise ErrorProcesamientoArchivo(mensaje_error_extraccion_final, ruta_archivo=ruta_archivo_entrada)


class ProcesadorArchivosMarkdown(ProcesadorArchivoInterfaz):
    """Procesador especializado para archivos Markdown (ej. .md, .markdown)."""
    EXTENSIONES_ARCHIVOS_SOPORTADAS = [".md", ".markdown"]

    def extraer_texto_de_archivo(self, ruta_archivo_entrada: Path) -> str:
        registrador.info(f"Extrayendo texto del archivo Markdown: '{ruta_archivo_entrada}'.")
        try:
            # Los archivos Markdown suelen estar codificados en UTF-8.
            with open(ruta_archivo_entrada, "r", encoding="utf-8") as archivo_md:
                texto_extraido_md = archivo_md.read()
            registrador.info(f"Texto extraído correctamente del archivo Markdown '{ruta_archivo_entrada}'.")
            return texto_extraido_md
        except IOError as e_error_io_md:
            mensaje_error_md = f"Error de E/S al leer el archivo Markdown '{ruta_archivo_entrada}': {e_error_io_md}"
            registrador.error(mensaje_error_md)
            raise ErrorProcesamientoArchivo(mensaje_error_md, e_error_io_md, ruta_archivo=ruta_archivo_entrada) from e_error_io_md
        except Exception as e_error_inesperado_md: # Capturar otros posibles errores
            mensaje_error_inesperado_md = f"Error inesperado al extraer texto del archivo Markdown '{ruta_archivo_entrada}': {e_error_inesperado_md}"
            registrador.exception(mensaje_error_inesperado_md) # Usar exception para incluir traceback completo en logs
            raise ErrorProcesamientoArchivo(mensaje_error_inesperado_md, e_error_inesperado_md, ruta_archivo=ruta_archivo_entrada) from e_error_inesperado_md


class ProcesadorArchivosPDF(ProcesadorArchivoInterfaz):
    """
    Procesador especializado para archivos PDF. Intenta extraer texto directamente si es posible (no implementado aquí),
    y como fallback utiliza OCR (Tesseract) convirtiendo las páginas del PDF a imágenes.
    """
    EXTENSIONES_ARCHIVOS_SOPORTADAS = [".pdf"]

    def __init__(self, lenguaje_ocr_predeterminado: str = "spa+eng"): # Español e Inglés por defecto para OCR
        if pytesseract is None or convert_from_path is None:
            mensaje_advertencia_dependencia = "Dependencias 'pytesseract' y/o 'pdf2image' no están instaladas. El procesamiento de PDF con OCR no estará completamente funcional."
            registrador.warning(mensaje_advertencia_dependencia)
            # No se lanza ErrorDependenciaFaltante aquí para permitir que la aplicación inicie.
            # El error se lanzará si se intenta usar `extraer_texto_de_archivo`.
        self.lenguaje_ocr = lenguaje_ocr_predeterminado
        registrador.debug(f"ProcesadorArchivosPDF inicializado. Lenguaje OCR predeterminado: '{self.lenguaje_ocr}'.")


    def extraer_texto_de_archivo(self, ruta_archivo_entrada: Path) -> str:
        if pytesseract is None or convert_from_path is None:
            # Comprobación de dependencias en tiempo de ejecución del método.
            mensaje_error_dependencia_pdf = "Faltan dependencias cruciales ('pytesseract' y/o 'pdf2image') para procesar archivos PDF."
            registrador.error(mensaje_error_dependencia_pdf)
            raise ErrorDependenciaFaltante(mensaje_error_dependencia_pdf, ruta_archivo=ruta_archivo_entrada)

        registrador.info(f"Iniciando extracción de texto del PDF: '{ruta_archivo_entrada}' mediante OCR.")
        lista_textos_por_pagina: List[str] = []
        try:
            # Convertir páginas del PDF a una lista de imágenes (objetos PIL.Image)
            registrador.debug(f"Convirtiendo PDF '{ruta_archivo_entrada}' a imágenes para OCR...")
            lista_imagenes_de_paginas_pdf = convert_from_path(ruta_archivo_entrada, timeout=60) # Timeout para la conversión

            if not lista_imagenes_de_paginas_pdf:
                registrador.warning(f"No se pudieron convertir páginas a imágenes desde el PDF '{ruta_archivo_entrada}'. El PDF podría estar vacío, corrupto o protegido contra extracción.")
                return "" # Devolver string vacío si no se obtuvieron imágenes

            registrador.info(f"PDF '{ruta_archivo_entrada}' convertido a {len(lista_imagenes_de_paginas_pdf)} imágenes. Procediendo con OCR en cada página...")
            for i, imagen_pagina_pdf_actual in enumerate(lista_imagenes_de_paginas_pdf):
                try:
                    registrador.debug(f"Procesando OCR para página {i + 1} del PDF '{ruta_archivo_entrada}'...")
                    # Extraer texto de la imagen de la página usando Tesseract OCR
                    texto_extraido_pagina_actual = pytesseract.image_to_string(imagen_pagina_pdf_actual, lang=self.lenguaje_ocr, timeout=30) # Timeout para OCR por página
                    if texto_extraido_pagina_actual and texto_extraido_pagina_actual.strip():
                        lista_textos_por_pagina.append(texto_extraido_pagina_actual.strip())
                        registrador.debug(f"Texto extraído de página {i + 1} (longitud: {len(texto_extraido_pagina_actual.strip())}).")
                    else:
                        registrador.debug(f"No se extrajo texto de la página {i + 1} del PDF '{ruta_archivo_entrada}' (página posiblemente vacía o sin texto detectable por OCR).")
                except pytesseract.TesseractError as error_tesseract_ocr: # Errores específicos de Tesseract
                    registrador.warning(f"Error de Tesseract OCR procesando página {i + 1} del PDF '{ruta_archivo_entrada}': {error_tesseract_ocr}. Se omitirá esta página.")
                except Exception as e_procesamiento_pagina_pdf: # Otros errores al procesar una imagen de página
                    registrador.warning(f"Error inesperado procesando la imagen de la página {i + 1} del PDF '{ruta_archivo_entrada}': {e_procesamiento_pagina_pdf}. Se omitirá esta página.")

            if not lista_textos_por_pagina and lista_imagenes_de_paginas_pdf: # Si se procesaron imágenes pero no se obtuvo texto
                 registrador.warning(f"No se extrajo texto de ninguna página del PDF '{ruta_archivo_entrada}' mediante OCR. El PDF podría no contener texto legible o el OCR falló consistentemente en todas las páginas.")

            texto_completo_extraido_pdf = "\n\n".join(filter(None, lista_textos_por_pagina)) # Unir textos de páginas, con doble salto de línea como separador
            registrador.info(f"Extracción de texto del PDF '{ruta_archivo_entrada}' completada. Páginas procesadas: {len(lista_imagenes_de_paginas_pdf)}, fragmentos de texto con contenido: {len(lista_textos_por_pagina)}.")
            return texto_completo_extraido_pdf

        except Exception as e_error_general_pdf: # Captura errores de convert_from_path o cualquier otro no previsto
            mensaje_error_pdf_general = f"No se pudo extraer texto del PDF '{ruta_archivo_entrada}' debido a un error general: {e_error_general_pdf}"
            registrador.exception(mensaje_error_pdf_general) # Loguear con traceback
            raise ErrorProcesamientoArchivo(mensaje_error_pdf_general, e_error_general_pdf, ruta_archivo=ruta_archivo_entrada) from e_error_general_pdf


class ProcesadorArchivosDocx(ProcesadorArchivoInterfaz):
    """Procesador especializado para archivos DOCX (formato de Microsoft Word)."""
    EXTENSIONES_ARCHIVOS_SOPORTADAS = [".docx"]

    def __init__(self):
        if docx is None: # Comprobar dependencia al inicializar
            registrador.warning("Dependencia 'python-docx' no está instalada. El procesador de DOCX no estará funcional.")
            # No lanzar error aquí para permitir que la aplicación inicie.

    def extraer_texto_de_archivo(self, ruta_archivo_entrada: Path) -> str:
        if docx is None: # Comprobar dependencia en tiempo de ejecución
            mensaje_error_dependencia_docx = "La dependencia 'python-docx' no está instalada. No se puede procesar el archivo DOCX."
            registrador.error(mensaje_error_dependencia_docx)
            raise ErrorDependenciaFaltante(mensaje_error_dependencia_docx, ruta_archivo=ruta_archivo_entrada)

        registrador.info(f"Extrayendo texto del archivo DOCX: '{ruta_archivo_entrada}'.")
        try:
            documento_word_abierto = docx.Document(str(ruta_archivo_entrada)) # python-docx espera un string con la ruta o un stream

            # Extraer texto de los párrafos del documento
            textos_parrafos_documento = [parrafo.text.strip() for parrafo in documento_word_abierto.paragraphs if parrafo.text and parrafo.text.strip()]

            # Extraer texto de las tablas del documento
            textos_celdas_tablas = []
            for tabla_doc in documento_word_abierto.tables:
                for fila_tabla in tabla_doc.rows:
                    for celda_tabla in fila_tabla.cells:
                        if celda_tabla.text and celda_tabla.text.strip():
                            textos_celdas_tablas.append(celda_tabla.text.strip())

            # Combinar todo el texto extraído de párrafos y tablas
            texto_completo_extraido_docx = "\n\n".join(textos_parrafos_documento + textos_celdas_tablas)
            registrador.info(f"Texto extraído correctamente del archivo DOCX '{ruta_archivo_entrada}'. Longitud total: {len(texto_completo_extraido_docx)} caracteres.")
            return texto_completo_extraido_docx
        except Exception as e_error_docx: # Capturar excepciones específicas de python-docx si se conocen, o genéricas
            mensaje_error_docx = f"No se pudo extraer texto del archivo DOCX '{ruta_archivo_entrada}': {e_error_docx}"
            registrador.exception(mensaje_error_docx)
            raise ErrorProcesamientoArchivo(mensaje_error_docx, e_error_docx, ruta_archivo=ruta_archivo_entrada) from e_error_docx


class ProcesadorArchivosPptx(ProcesadorArchivoInterfaz):
    """Procesador especializado para archivos PPTX (formato de Microsoft PowerPoint)."""
    EXTENSIONES_ARCHIVOS_SOPORTADAS = [".pptx"]

    def __init__(self):
        if Presentation is None: # Comprobar dependencia al inicializar
            registrador.warning("Dependencia 'python-pptx' no está instalada. El procesador de PPTX no estará funcional.")
            # No lanzar error aquí.

    def extraer_texto_de_archivo(self, ruta_archivo_entrada: Path) -> str:
        if Presentation is None: # Comprobar dependencia en tiempo de ejecución
            mensaje_error_dependencia_pptx = "La dependencia 'python-pptx' no está instalada. No se puede procesar el archivo PPTX."
            registrador.error(mensaje_error_dependencia_pptx)
            raise ErrorDependenciaFaltante(mensaje_error_dependencia_pptx, ruta_archivo=ruta_archivo_entrada)

        registrador.info(f"Extrayendo texto del archivo PPTX: '{ruta_archivo_entrada}'.")
        try:
            presentacion_powerpoint_abierta = Presentation(str(ruta_archivo_entrada)) # python-pptx espera un string con la ruta o un stream

            textos_extraidos_diapositivas = []
            for i, diapositiva_actual_ppt in enumerate(presentacion_powerpoint_abierta.slides):
                textos_formas_diapositiva_actual = []
                # Extraer texto de las formas (shapes) en cada diapositiva
                for forma_ppt_actual in diapositiva_actual_ppt.shapes:
                    if hasattr(forma_ppt_actual, "text_frame") and forma_ppt_actual.text_frame and forma_ppt_actual.text_frame.text:
                        textos_formas_diapositiva_actual.append(forma_ppt_actual.text_frame.text.strip())
                    elif hasattr(forma_ppt_actual, "text") and forma_ppt_actual.text: # Algunas formas simples pueden tener .text directamente
                        textos_formas_diapositiva_actual.append(forma_ppt_actual.text.strip())

                # Extraer texto de las notas de la diapositiva, si existen y tienen contenido
                if diapositiva_actual_ppt.has_notes_slide and \
                   diapositiva_actual_ppt.notes_slide and \
                   diapositiva_actual_ppt.notes_slide.notes_text_frame and \
                   diapositiva_actual_ppt.notes_slide.notes_text_frame.text:
                    texto_notas_diapositiva_actual = diapositiva_actual_ppt.notes_slide.notes_text_frame.text.strip()
                    if texto_notas_diapositiva_actual: # Solo añadir si hay contenido en las notas
                        textos_formas_diapositiva_actual.append(f"\n[Notas de Diapositiva {i + 1}]:\n{texto_notas_diapositiva_actual}")

                texto_consolidado_diapositiva_actual = "\n".join(filter(None, textos_formas_diapositiva_actual)) # Unir textos de formas y notas de la diapositiva
                if texto_consolidado_diapositiva_actual: # Solo añadir si la diapositiva tiene texto
                    textos_extraidos_diapositivas.append(f"[Contenido Diapositiva {i + 1}]:\n{texto_consolidado_diapositiva_actual}")

            texto_completo_extraido_pptx = "\n\n".join(filter(None, textos_extraidos_diapositivas)) # Unir textos de todas las diapositivas
            registrador.info(f"Texto extraído correctamente del archivo PPTX '{ruta_archivo_entrada}'. Longitud total: {len(texto_completo_extraido_pptx)}.")
            return texto_completo_extraido_pptx
        except Exception as e_error_pptx: # Capturar excepciones específicas de python-pptx si se conocen, o genéricas
            mensaje_error_pptx = f"No se pudo extraer texto del archivo PPTX '{ruta_archivo_entrada}': {e_error_pptx}"
            registrador.exception(mensaje_error_pptx)
            raise ErrorProcesamientoArchivo(mensaje_error_pptx, e_error_pptx, ruta_archivo=ruta_archivo_entrada) from e_error_pptx

# --- Gestor Principal de Procesadores de Archivos ---

class GestorMaestroDeProcesadoresArchivos:
    """
    Clase central que gestiona una colección de procesadores de archivos.
    Delega la tarea de procesamiento de un archivo al procesador adecuado
    según la extensión del archivo.
    """
    def __init__(self):
        self.mapeo_procesadores_por_extension: Dict[str, ProcesadorArchivoInterfaz] = {}
        self._registrar_procesadores_disponibles_por_defecto() # Registrar procesadores al inicializar
        registrador.info("GestorMaestroDeProcesadoresArchivos inicializado con procesadores de archivo por defecto.")

    def _registrar_procesadores_disponibles_por_defecto(self):
        """
        Registra instancias de los procesadores de archivo por defecto que están disponibles
        (es decir, cuyas dependencias de software están cumplidas).
        """
        self.intentar_registrar_procesador(ProcesadorArchivosTextoPlano())
        self.intentar_registrar_procesador(ProcesadorArchivosMarkdown())

        # Registrar procesador de PDF solo si las dependencias están presentes
        if pytesseract and convert_from_path:
            self.intentar_registrar_procesador(ProcesadorArchivosPDF())
        else:
            registrador.warning("Procesador de PDF (ProcesadorArchivosPDF) no será registrado debido a que faltan las dependencias 'pytesseract' y/o 'pdf2image'.")

        # Registrar procesador de DOCX solo si la dependencia está presente
        if docx:
            self.intentar_registrar_procesador(ProcesadorArchivosDocx())
        else:
            registrador.warning("Procesador de DOCX (ProcesadorArchivosDocx) no será registrado porque falta la dependencia 'python-docx'.")

        # Registrar procesador de PPTX solo si la dependencia está presente
        if Presentation: # Clase Presentation de la biblioteca python-pptx
            self.intentar_registrar_procesador(ProcesadorArchivosPptx())
        else:
            registrador.warning("Procesador de PPTX (ProcesadorArchivosPptx) no será registrado porque falta la dependencia 'python-pptx'.")

    def intentar_registrar_procesador(self, procesador_para_registrar: ProcesadorArchivoInterfaz): # Parámetro renombrado
        """
        Registra un procesador de archivos para cada una de las extensiones que este soporta.
        Si una extensión ya tiene un procesador registrado, el nuevo procesador lo sobrescribirá.

        Args:
            procesador_para_registrar: Instancia del procesador a registrar.
        """
        if not isinstance(procesador_para_registrar, ProcesadorArchivoInterfaz): # Verificar tipo
            registrador.error(f"Intento de registrar un objeto que no es una instancia válida de ProcesadorArchivoInterfaz: {type(procesador_para_registrar)}")
            return

        for extension_soportada_actual in procesador_para_registrar.EXTENSIONES_ARCHIVOS_SOPORTADAS:
            extension_normalizada_actual = extension_soportada_actual.lower() # Asegurar minúsculas
            self.mapeo_procesadores_por_extension[extension_normalizada_actual] = procesador_para_registrar
            registrador.info(f"Procesador para extensión '{extension_normalizada_actual}' registrado: {type(procesador_para_registrar).__name__}")

    def procesar_archivo_segun_tipo(self, ruta_archivo_entrada_a_procesar: Path) -> Optional[str]: # Parámetro renombrado
        """
        Procesa un archivo utilizando el procesador adecuado según su extensión.

        Args:
            ruta_archivo_entrada_a_procesar: Objeto Path que apunta al archivo a procesar.

        Returns:
            El texto extraído como un string, o None si el archivo no existe,
            no hay un procesador registrado para su extensión, o si ocurre un error
            durante el procesamiento y el procesador específico no maneja la excepción.
        """
        if not isinstance(ruta_archivo_entrada_a_procesar, Path): # Verificar tipo de entrada
            try:
                # Intentar convertir a Path si se recibió un string, por flexibilidad.
                ruta_archivo_entrada_a_procesar = Path(ruta_archivo_entrada_a_procesar)
            except TypeError: # Si no se puede convertir a Path (ej. es None o un tipo incompatible)
                registrador.error(f"La ruta del archivo proporcionada es inválida o de tipo incorrecto: '{ruta_archivo_entrada_a_procesar}' (tipo: {type(ruta_archivo_entrada_a_procesar)}).")
                return None # No se puede procesar

        registrador.info(f"Solicitud para procesar archivo: '{ruta_archivo_entrada_a_procesar}'.")
        if not ruta_archivo_entrada_a_procesar.is_file(): # Verificar si el archivo existe y es un archivo
            registrador.error(f"El archivo especificado '{ruta_archivo_entrada_a_procesar}' no existe o no es un archivo válido.")
            return None # Archivo no encontrado o no es un archivo

        extension_archivo_a_procesar = ruta_archivo_entrada_a_procesar.suffix.lower() # Obtener extensión en minúsculas
        procesador_seleccionado_para_extension = self.mapeo_procesadores_por_extension.get(extension_archivo_a_procesar)

        if procesador_seleccionado_para_extension:
            nombre_procesador_seleccionado = type(procesador_seleccionado_para_extension).__name__
            registrador.info(f"Procesando archivo '{ruta_archivo_entrada_a_procesar}' con el procesador: {nombre_procesador_seleccionado}.")
            try:
                texto_extraido_del_archivo = procesador_seleccionado_para_extension.extraer_texto_de_archivo(ruta_archivo_entrada_a_procesar)
                registrador.info(f"Procesamiento de '{ruta_archivo_entrada_a_procesar}' con '{nombre_procesador_seleccionado}' finalizado. Longitud del texto extraído: {len(texto_extraido_del_archivo) if texto_extraido_del_archivo is not None else 'N/A'}.")
                return texto_extraido_del_archivo
            except ErrorDependenciaFaltante as e_error_dependencia: # Capturar error de dependencia específica si es lanzado por el procesador
                registrador.error(f"Error de dependencia faltante al procesar '{ruta_archivo_entrada_a_procesar}' con {nombre_procesador_seleccionado}: {e_error_dependencia}")
                # No relanzar, simplemente devolver None para indicar fallo de procesamiento.
                return None
            except ErrorProcesamientoArchivo as e_error_procesamiento_archivo: # Otros errores de procesamiento definidos
                registrador.error(f"Error específico de procesamiento al procesar '{ruta_archivo_entrada_a_procesar}' con {nombre_procesador_seleccionado}: {e_error_procesamiento_archivo}")
                return None # Devolver None para indicar fallo
            except Exception as e_error_inesperado_procesamiento: # Errores completamente inesperados no capturados por el procesador
                 registrador.exception(f"Error inesperado y no capturado durante el procesamiento de '{ruta_archivo_entrada_a_procesar}' con {nombre_procesador_seleccionado}: {e_error_inesperado_procesamiento}")
                 return None # Devolver None para indicar fallo
        else: # No se encontró un procesador para la extensión
            mensaje_tipo_no_soportado = f"No se encontró un procesador adecuado para la extensión '{extension_archivo_a_procesar}' del archivo '{ruta_archivo_entrada_a_procesar}'. Archivo no será procesado."
            registrador.warning(mensaje_tipo_no_soportado)
            # Considerar si se debe lanzar ErrorTipoArchivoNoSoportado aquí o simplemente devolver None.
            # Devolver None es más suave si se esperan archivos no soportados.
            # Si se espera que todos los archivos sean soportados, lanzar una excepción sería mejor.
            # Por ahora, se devuelve None, y el llamador puede verificar si el resultado es None.
            return None
[end of entrenai_refactor/nucleo/archivos/procesador_archivos.py]
