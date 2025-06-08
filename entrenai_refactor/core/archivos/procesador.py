import logging
from pathlib import Path
from typing import Optional, Dict, List

import docx
import pytesseract
from pdf2image import convert_from_path
from pptx import Presentation

# Configuración de logger en español
logger = logging.getLogger("procesador_archivos")
if not logger.handlers:
    handler = logging.StreamHandler()
    formatter = logging.Formatter("%(asctime)s - %(levelname)s - %(message)s")
    handler.setFormatter(formatter)
    logger.addHandler(handler)
logger.setLevel(logging.INFO)

class ErrorProcesamientoArchivo(Exception):
    pass

class ProcesadorBase:
    EXTENSIONES_SOPORTADAS: List[str] = []

    def extraer_texto(self, ruta_archivo: Path) -> str:
        raise NotImplementedError

    def puede_procesar(self, ruta_archivo: Path) -> bool:
        return ruta_archivo.suffix.lower() in self.EXTENSIONES_SOPORTADAS

class ProcesadorTxt(ProcesadorBase):
    EXTENSIONES_SOPORTADAS = [".txt"]

    def extraer_texto(self, ruta_archivo: Path) -> str:
        for codificacion in ["utf-8", "latin-1", "iso-8859-1", "cp1252"]:
            try:
                with open(ruta_archivo, "r", encoding=codificacion) as f:
                    texto = f.read()
                    logger.info(f"Texto extraído correctamente de {ruta_archivo}")
                    return texto
            except Exception:
                continue
        logger.error(f"No se pudo extraer texto del archivo TXT {ruta_archivo}")
        raise ErrorProcesamientoArchivo(f"No se pudo extraer texto del archivo TXT {ruta_archivo}")

class ProcesadorMarkdown(ProcesadorBase):
    EXTENSIONES_SOPORTADAS = [".md", ".markdown"]

    def extraer_texto(self, ruta_archivo: Path) -> str:
        try:
            with open(ruta_archivo, "r", encoding="utf-8") as f:
                texto = f.read()
                logger.info(f"Texto extraído correctamente de {ruta_archivo}")
                return texto
        except Exception as e:
            logger.error(f"No se pudo extraer texto del archivo Markdown {ruta_archivo}: {e}")
            raise ErrorProcesamientoArchivo(f"No se pudo extraer texto del archivo Markdown {ruta_archivo}")

class ProcesadorPdf(ProcesadorBase):
    EXTENSIONES_SOPORTADAS = [".pdf"]

    def extraer_texto(self, ruta_archivo: Path) -> str:
        partes = []
        try:
            imagenes = convert_from_path(ruta_archivo)
            for imagen in imagenes:
                partes.append(pytesseract.image_to_string(imagen, lang="spa+eng"))
            logger.info(f"Texto extraído correctamente de {ruta_archivo}")
            return "\n\n".join(partes)
        except Exception as e:
            logger.error(f"No se pudo extraer texto del PDF {ruta_archivo}: {e}")
            raise ErrorProcesamientoArchivo(f"No se pudo extraer texto del PDF {ruta_archivo}")

class ProcesadorDocx(ProcesadorBase):
    EXTENSIONES_SOPORTADAS = [".docx"]

    def extraer_texto(self, ruta_archivo: Path) -> str:
        try:
            documento = docx.Document(str(ruta_archivo))
            texto = []
            for parrafo in documento.paragraphs:
                texto.append(parrafo.text)
            for tabla in documento.tables:
                for fila in tabla.rows:
                    for celda in fila.cells:
                        texto.append(celda.text)
            logger.info(f"Texto extraído correctamente de {ruta_archivo}")
            return "\n\n".join(texto)
        except Exception as e:
            logger.error(f"No se pudo extraer texto del archivo DOCX {ruta_archivo}: {e}")
            raise ErrorProcesamientoArchivo(f"No se pudo extraer texto del archivo DOCX {ruta_archivo}")

class ProcesadorPptx(ProcesadorBase):
    EXTENSIONES_SOPORTADAS = [".pptx"]

    def extraer_texto(self, ruta_archivo: Path) -> str:
        try:
            presentacion = Presentation(str(ruta_archivo))
            texto = []
            for diapositiva in presentacion.slides:
                for forma in diapositiva.shapes:
                    if hasattr(forma, "text_frame") and forma.text_frame:
                        for parrafo in forma.text_frame.paragraphs:
                            for run in parrafo.runs:
                                texto.append(run.text)
                    elif hasattr(forma, "text"):
                        texto.append(forma.text)
                if diapositiva.has_notes_slide and diapositiva.notes_slide.notes_text_frame:
                    texto.append(diapositiva.notes_slide.notes_text_frame.text)
            logger.info(f"Texto extraído correctamente de {ruta_archivo}")
            return "\n\n".join(filter(None, texto))
        except Exception as e:
            logger.error(f"No se pudo extraer texto del archivo PPTX {ruta_archivo}: {e}")
            raise ErrorProcesamientoArchivo(f"No se pudo extraer texto del archivo PPTX {ruta_archivo}")

class ProcesadorArchivos:
    def __init__(self):
        self.procesadores: Dict[str, ProcesadorBase] = {}
        self._registrar_procesadores()

    def _registrar_procesadores(self):
        self.registrar_procesador(ProcesadorTxt())
        self.registrar_procesador(ProcesadorMarkdown())
        self.registrar_procesador(ProcesadorPdf())
        self.registrar_procesador(ProcesadorDocx())
        self.registrar_procesador(ProcesadorPptx())

    def registrar_procesador(self, procesador: ProcesadorBase):
        for ext in procesador.EXTENSIONES_SOPORTADAS:
            self.procesadores[ext.lower()] = procesador

    def procesar(self, ruta_archivo: str) -> Optional[str]:
        path = Path(ruta_archivo)
        if not path.is_file():
            logger.warning(f"El archivo {ruta_archivo} no existe.")
            return None
        ext = path.suffix.lower()
        procesador = self.procesadores.get(ext)
        if procesador:
            try:
                return procesador.extraer_texto(path)
            except ErrorProcesamientoArchivo:
                return None
        logger.warning(f"No hay procesador para la extensión {ext}")
        return None 